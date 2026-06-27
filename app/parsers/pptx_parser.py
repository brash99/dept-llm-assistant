from pathlib import Path

from pptx import Presentation

from app.parser import Parser
from app.knowledge import document_from_text


class PPTXParser(Parser):

    name = "PPTXParser"
    supported_suffixes = {".pptx"}

    def parse(self, path, root_path):

        path = Path(path)

        prs = Presentation(path)

        slides = []
        slide_lengths = []

        for i, slide in enumerate(prs.slides, start=1):

            text = [f"--- Slide {i} ---"]

            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    s = shape.text.strip()
                    if s:
                        text.append(s)

            slide_text = "\n".join(text)

            slides.append(slide_text)
            slide_lengths.append(len(slide_text))

        full_text = "\n\n".join(slides).strip()

        cp = prs.core_properties

        metadata = {
            "num_slides": len(prs.slides),
            "slide_lengths": slide_lengths,
            "text_length": len(full_text),
            "quality": "good" if len(full_text) >= 100 else "poor",
            "is_empty": len(full_text) == 0,
            "core_properties": {
                "author": cp.author,
                "title": cp.title,
                "subject": cp.subject,
                "keywords": cp.keywords,
                "comments": cp.comments,
                "category": cp.category,
                "created": cp.created.isoformat() if cp.created else None,
                "modified": cp.modified.isoformat() if cp.modified else None,
                "last_modified_by": cp.last_modified_by,
            },
        }

        title = cp.title or path.stem

        return document_from_text(
            source_path=path,
            root_path=root_path,
            text=full_text,
            parser=self.name,
            title=title,
            metadata=metadata,
        )
